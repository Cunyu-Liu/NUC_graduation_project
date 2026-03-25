#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建中期答辩PPT - 根据中期检查表详细内容
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色主题 - 学术蓝色调
PRIMARY_COLOR = RGBColor(0x1E, 0x3A, 0x5F)  # 深蓝
ACCENT_COLOR = RGBColor(0x2E, 0x86, 0xAB)   # 中蓝
LIGHT_COLOR = RGBColor(0xA8, 0xD0, 0xE2)    # 浅蓝
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
SUCCESS_COLOR = RGBColor(0x28, 0xA7, 0x45)  # 绿色
WARNING_COLOR = RGBColor(0xFF, 0xA5, 0x00)  # 橙色

def add_title_slide(prs, title, subtitle, info):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                    prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # 添加学校名称
    school_box = slide.shapes.add_textbox(Inches(0), Inches(1.5), prs.slide_width, Inches(1))
    tf = school_box.text_frame
    p = tf.paragraphs[0]
    p.text = "中北大学软件学院"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 添加答辩类型标签
    type_box = slide.shapes.add_textbox(Inches(0), Inches(2.3), prs.slide_width, Inches(0.6))
    tf = type_box.text_frame
    p = tf.paragraphs[0]
    p.text = "毕业设计中期答辩"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = LIGHT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # 添加个人信息
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.333), Inches(2))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    info_lines = info.split('\n')
    for i, line in enumerate(info_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)
    
    return slide

def add_content_slide(prs, title, content_items, subtitle=None):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 添加副标题（如果有）
    y_start = 1.5
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_COLOR
        p.font.italic = True
        y_start = 1.9
    
    # 添加内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_start), Inches(12), Inches(5.8 - y_start))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(item, tuple):
            # (标题, 内容) 格式
            p.text = f"● {item[0]}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
            p.space_before = Pt(10)
            p.space_after = Pt(3)
            
            p2 = tf.add_paragraph()
            p2.text = f"   {item[1]}"
            p2.font.size = Pt(15)
            p2.font.color.rgb = DARK_TEXT
            p2.space_after = Pt(6)
        else:
            p.text = f"● {item}"
            p.font.size = Pt(17)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(8)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
    
    # 添加分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.5), 
                                   Inches(0.02), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    return slide

def add_timeline_slide(prs, title, timeline_items):
    """添加时间线/进度页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 添加时间线内容
    y_pos = 1.6
    for i, (phase, desc, status) in enumerate(timeline_items):
        # 状态颜色
        if status == "completed":
            status_color = SUCCESS_COLOR
            status_text = "✓ 已完成"
        elif status == "in_progress":
            status_color = WARNING_COLOR
            status_text = "▶ 进行中"
        else:
            status_color = RGBColor(0x99, 0x99, 0x99)
            status_text = "○ 待开始"
        
        # 阶段标题
        phase_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(2.8), Inches(0.5))
        tf = phase_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(3.7), Inches(y_pos), Inches(7), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_TEXT
        
        # 状态标签
        status_box = slide.shapes.add_textbox(Inches(11), Inches(y_pos), Inches(1.8), Inches(0.5))
        tf = status_box.text_frame
        p = tf.paragraphs[0]
        p.text = status_text
        p.font.size = Pt(14)
        p.font.color.rgb = status_color
        p.font.bold = True
        
        y_pos += 0.9
    
    return slide

def add_problem_slide(prs, title, problems, solutions):
    """添加问题与对策页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 左侧：存在的问题
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "存在的问题"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)  # 红色
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(problems):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)
    
    # 右侧：解决方案
    right_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "解决方案"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_COLOR
    
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(solutions):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)
    
    # 添加分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.5), 
                                   Inches(0.02), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    return slide

# ==================== 创建PPT内容 ====================

# 第1页：标题页
add_title_slide(prs,
    "科研文献摘要提取系统的设计与实现",
    "基于大语言模型的螺旋式知识积累与代码生成平台",
    "学  号：2213041523\n姓  名：刘存宇\n班  级：22130415\n指导教师：王斌"
)

# 第2页：研究背景与意义
add_content_slide(prs, "研究背景与意义", [
    ("科研文献爆炸式增长", "全球每年发表论文超数百万篇，传统人工阅读模式效率低下，科研人员面临严重的信息过载困境"),
    ("大模型技术突破", "GPT-4、GLM-4等大语言模型在文本理解、摘要生成方面展现出强大能力，为科研辅助提供新可能"),
    ("螺旋式知识积累", "构建从文献分析到代码生成的完整闭环，实现知识沉淀与复用，形成持续迭代的研究模式"),
    "提升文献处理效率：将小时级阅读压缩至分钟级分析",
    "精准提取关键信息：智能识别12类核心要点（创新点、方法、实验等）",
    "辅助学术前沿发现：通过主题聚类识别研究热点与空白"
])

# 第3页：系统架构设计
add_two_column_slide(prs, "系统架构与技术栈",
    "后端架构", [
        "Flask 3.0 - Web服务框架",
        "PostgreSQL - 主数据库(关系型)",
        "Milvus - 向量数据库(Docker部署)",
        "SQLAlchemy 2.0 - ORM框架",
        "LangChain - LLM链式编排",
        "GLM-4 API - 智谱AI大模型",
        "Redis - 缓存层(可选)",
        "Socket.IO - WebSocket实时通信"
    ],
    "前端架构", [
        "Vue 3 - 响应式前端框架",
        "Element Plus - UI组件库",
        "D3.js - 力导向图谱可视化",
        "Monaco Editor - 代码编辑器",
        "Axios - HTTP请求封装",
        "Socket.IO - 实时进度推送",
        "Markdown-it - Markdown渲染",
        "KaTeX - LaTeX公式渲染"
    ]
)

# 第4页：核心功能模块 - 已完成（基于中期检查表详细内容）
add_content_slide(prs, "已完成工作 - 核心功能开发", [
    ("PDF智能解析模块", "实现中英文PDF文档的高精度解析，支持文本、表格、图片提取，为元数据自动识别提供基础"),
    ("AI摘要生成模块", "基于GLM-4 API实现博士级学术摘要自动生成，采用专业提示词模板确保学术规范"),
    ("12类要点提取", "智能提取论文中的创新点、研究方法、实验设计、关键结论等12类核心科研要素"),
    ("研究空白挖掘", "自动识别5种类型的研究空白，并结合重要性+难度进行智能优先级排序")
], "核心数据处理功能")

# 第5页：知识图谱与可视化
add_content_slide(prs, "已完成工作 - 知识图谱与可视化", [
    ("D3.js力导向布局", "基于D3.js实现交互式知识图谱可视化，节点可拖拽、缩放，支持多种关系类型展示"),
    ("论文关系网络", "自动构建论文关联网络，支持5种关系类型：引用、相似、递进、对比、互补"),
    ("交互功能", "实现节点点击详情、关系筛选、图谱缩放、拖拽定位等丰富的用户交互体验"),
    ("动态更新", "支持新增论文后自动更新图谱结构，实时反映知识库变化")
], "知识图谱可视化模块")

# 第6页：AI智能助手功能
add_two_column_slide(prs, "已完成工作 - AI智能助手功能",
    "Kimi风格AI聊天", [
        "流式输出，实时响应",
        "Markdown语法完整渲染",
        "LaTeX数学公式支持",
        "代码块语法高亮",
        "论文引用自动标注"
    ],
    "RAG与链式工作流", [
        "基于RAG的论文库问答",
        "向量检索相似文献",
        "LangChain SequentialChain",
        "多步骤分析流程编排",
        "预设模板一键分析"
    ]
)

# 第7页：向量聚类与代码生成
add_content_slide(prs, "已完成工作 - 向量聚类与代码生成", [
    ("Milvus向量数据库", "集成Milvus向量数据库，基于BGE-large模型生成语义嵌入，实现论文语义相似度计算"),
    ("语义聚类分析", "支持基于语义相似度的论文自动聚类，发现研究主题群组，推荐相似论文"),
    ("6种代码生成策略", "方法改进、新方法提出、数据集创建、实验设计、模型实现、算法优化"),
    ("代码编辑器集成", "集成Monaco Editor提供专业代码编辑体验，支持代码版本历史管理")
], "智能分析与生成模块")

# 第8页：性能优化成果
add_two_column_slide(prs, "性能优化成果",
    "异步并发优化", [
        "异步工作流引擎重构",
        "支持100篇论文并发处理",
        "分析速度提升6倍",
        "Semaphore并发控制",
        "WebSocket实时进度推送"
    ],
    "数据库与缓存优化", [
        "创建30+数据库索引",
        "查询速度提升10-50倍",
        "Redis缓存层（命中率>80%）",
        "Gzip响应压缩（减少60%）",
        "连接池优化管理"
    ]
)

# 第9页：存在的问题
add_problem_slide(prs, "目前存在的问题与对策",
    [
        "系统测试覆盖度有待提升，部分边界情况处理不够完善",
        "大模型API调用成本较高，需要优化调用策略",
        "知识图谱在大规模数据下的渲染性能有待优化",
        "部分复杂PDF格式的解析准确率需要进一步提升"
    ],
    [
        "完善单元测试和集成测试，增加异常处理机制",
        "实现缓存机制和批量调用优化，降低成本",
        "引入图谱分层渲染和虚拟滚动技术优化性能",
        "集成多PDF解析引擎，提升复杂格式支持"
    ]
)

# 第10页：下一步工作计划
add_timeline_slide(prs, "下一步工作计划", [
    ("系统测试与优化", "完善单元测试、集成测试，进行压力测试和性能调优", "in_progress"),
    ("论文撰写", "完成毕业设计说明书初稿，完善系统使用说明文档", "pending"),
    ("系统部署", "完成Docker容器化部署配置，编写部署文档和用户手册", "pending"),
    ("答辩准备", "制作答辩PPT，准备系统演示环境，进行答辩演练", "pending"),
    ("最终验收", "提交所有材料，参加毕业答辩，根据反馈完善", "pending")
])

# 第11页：总结与展望
add_content_slide(prs, "总结与展望", [
    ("当前成果", "已完成系统核心功能开发，包括PDF解析、AI摘要、要点提取、知识图谱、代码生成等模块"),
    ("技术亮点", "螺旋式知识积累架构、100篇并发处理、Kimi风格AI聊天、LangChain链式工作流、Milvus向量聚类"),
    ("应用价值", "为科研人员提供从文献分析到代码生成的完整闭环支持，显著提升科研效率"),
    ("后续规划", "完善系统测试，完成论文撰写，做好答辩准备，确保顺利毕业验收"),
    "",
    "感谢各位老师的指导！"
])

# 保存PPT
output_path = "/Users/liucunyu/Documents/all_code/NUC_graduation_project/output/2213041523-刘存宇-中期答辩PPT.pptx"
prs.save(output_path)
print(f"✅ 中期答辩PPT已生成！")
print(f"📁 保存路径: {output_path}")
print(f"📊 共包含 {len(prs.slides)} 页幻灯片")
