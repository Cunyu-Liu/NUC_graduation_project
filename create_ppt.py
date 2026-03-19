#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建中期答辩PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色主题 - 学术蓝色调
PRIMARY_COLOR = RGBColor(0x1E, 0x3A, 0x5F)  # 深蓝
ACCENT_COLOR = RGBColor(0x2E, 0x86, 0xAB)   # 中蓝
LIGHT_COLOR = RGBColor(0xA8, 0xD0, 0xE2)    # 浅蓝
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

def add_title_slide(prs, title, subtitle, info):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                    prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # 添加学校名称
    school_box = slide.shapes.add_textbox(Inches(0), Inches(1.5), prs.slide_width, Inches(1))
    tf = school_box.text_frame
    p = tf.paragraphs[0]
    p.text = "中北大学软件学院"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # 添加个人信息
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.333), Inches(2))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    info_lines = info.split('\n')
    for i, line in enumerate(info_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)
    
    return slide

def add_content_slide(prs, title, content_items):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 添加内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(item, tuple):
            # (标题, 内容) 格式
            p.text = f"● {item[0]}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
            p.space_before = Pt(12)
            p.space_after = Pt(4)
            
            p2 = tf.add_paragraph()
            p2.text = f"   {item[1]}"
            p2.font.size = Pt(16)
            p2.font.color.rgb = DARK_TEXT
            p2.space_after = Pt(8)
        else:
            p.text = f"● {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(10)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
    
    # 添加分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.5), 
                                   Inches(0.02), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    return slide

def add_timeline_slide(prs, title, timeline_items):
    """添加时间线/进度页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 添加时间线内容
    y_pos = 1.6
    for i, (phase, desc, status) in enumerate(timeline_items):
        # 状态颜色
        if status == "completed":
            status_color = RGBColor(0x28, 0xA7, 0x45)  # 绿色
            status_text = "✓ 已完成"
        elif status == "in_progress":
            status_color = RGBColor(0xFF, 0xA5, 0x00)  # 橙色
            status_text = "▶ 进行中"
        else:
            status_color = RGBColor(0x99, 0x99, 0x99)  # 灰色
            status_text = "○ 待开始"
        
        # 阶段标题
        phase_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(3), Inches(0.5))
        tf = phase_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(4), Inches(y_pos), Inches(7), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        
        # 状态标签
        status_box = slide.shapes.add_textbox(Inches(11), Inches(y_pos), Inches(1.8), Inches(0.5))
        tf = status_box.text_frame
        p = tf.paragraphs[0]
        p.text = status_text
        p.font.size = Pt(14)
        p.font.color.rgb = status_color
        p.font.bold = True
        
        y_pos += 0.9
    
    return slide

# ==================== 创建PPT内容 ====================

# 第1页：标题页
add_title_slide(prs,
    "科研文献摘要提取系统的设计与实现",
    "基于大语言模型的螺旋式知识积累与代码生成平台",
    "学  号：2213041523\n姓  名：刘存宇\n班  级：22130415\n指导教师：王斌"
)

# 第2页：研究背景与意义
add_content_slide(prs, "研究背景与意义", [
    ("科研文献爆炸", "全球每年发表论文超数百万篇，传统阅读模式效率低下，科研人员面临信息过载困境"),
    ("大模型技术机遇", "GPT-4、GLM-4等大语言模型在文本理解、摘要生成方面取得突破性进展"),
    ("研究意义", "构建全流程科研辅助系统，实现从文献分析到代码生成的完整闭环"),
    "提升文献处理效率：将小时级阅读压缩至分钟级",
    "精准获取关键信息：智能提取创新点、方法、实验结论等12类要点",
    "辅助学术前沿发现：主题聚类识别研究热点与空白"
])

# 第3页：系统架构设计
add_two_column_slide(prs, "系统架构与技术栈",
    "后端架构", [
        "Flask 3.0 - Web框架",
        "PostgreSQL - 主数据库",
        "Milvus - 向量数据库",
        "SQLAlchemy 2.0 - ORM",
        "LangChain - LLM编排",
        "GLM-4 API - 大语言模型",
        "Redis - 缓存层"
    ],
    "前端架构", [
        "Vue 3 - 前端框架",
        "Element Plus - UI组件",
        "D3.js - 知识图谱可视化",
        "Monaco Editor - 代码编辑",
        "Socket.IO - 实时通信",
        "Axios - HTTP客户端"
    ]
)

# 第4页：已完成工作
add_content_slide(prs, "已完成的工作", [
    ("核心功能开发", "完成PDF智能解析、AI摘要生成、12类要点提取功能"),
    ("知识图谱模块", "基于D3.js实现力导向布局可视化，自动发现论文关联关系"),
    ("代码生成引擎", "支持6种生成策略：方法改进、新方法、数据集构建等"),
    ("AI聊天系统", "实现Kimi风格流式对话，支持RAG检索增强生成"),
    ("链式工作流", "基于LangChain SequentialChain的多步骤分析流程"),
    ("向量聚类", "集成Milvus向量数据库，实现语义相似度论文聚类"),
    ("性能优化", "异步并发支持100篇论文，查询速度提升10-50倍")
])

# 第5页：下一步计划
add_timeline_slide(prs, "下一步工作计划", [
    ("系统测试与优化", "完善单元测试、集成测试，修复潜在Bug", "in_progress"),
    ("论文撰写", "完成毕业设计说明书撰写", "pending"),
    ("系统部署", "Docker容器化部署，编写部署文档", "pending"),
    ("答辩准备", "制作答辩PPT，准备演示环境", "pending"),
    ("最终验收", "系统演示与答辩", "pending")
])

# 保存PPT
output_path = "/Users/liucunyu/Documents/all_code/NUC_graduation_project/2213041523-刘存宇-中期答辩PPT.pptx"
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
